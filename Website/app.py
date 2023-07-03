import streamlit as st
import openai
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
from langchain.chat_models import ChatOpenAI
from langchain.chains import LLMChain
from langchain.prompts.chat import (
    ChatPromptTemplate,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
)
from sentence_transformers import SentenceTransformer
import scipy.spatial
from docx import Document
from pptx import Presentation
import time

openai.api_key = st.secrets["OPENAI_API_KEY"]
embeddings = OpenAIEmbeddings()


def check_answer(question, user_answer, summary, embeddings):
    """
    Uses OpenAI to generate an answer to the selected question based on the summary.
    Compares the user's answer with the generated answer to determine if they convey the same meaning.
    """

    # Return False if the answer is empty, but also provide the correct answer
    if not user_answer.strip():
        # Generate answer from OpenAI model
        db = FAISS.from_texts([summary], embeddings)
        model_answer, _ = get_response_from_query(
            db, question, summary_needed=False)
        return False, model_answer

    # Generate answer from OpenAI model
    db = FAISS.from_texts([summary], embeddings)
    model_answer, _ = get_response_from_query(
        db, question, summary_needed=False)

    # Use SentenceTransformer to get sentence embeddings
    model = SentenceTransformer('all-MiniLM-L6-v2')
    user_answer_emb = model.encode([user_answer])[0]
    model_answer_emb = model.encode([model_answer])[0]

    # Calculate cosine similarity
    cosine_sim = 1 - \
        scipy.spatial.distance.cosine(user_answer_emb, model_answer_emb)

    # If the cosine similarity is above a threshold, consider it as correct
    if cosine_sim > 0.7:
        return True, ""
    else:
        return False, model_answer


def read_docx(file):
    doc = Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return ' '.join(full_text)


def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def generate_questions(summary, difficulty):
    """
    This function takes a summary and difficulty level as input and generates 5 questions using the OpenAI model.
    """
    if difficulty == 1:
        temperature = 0.1
        prompt = f"I have read a text about the following topic: {summary}. Please generate 5 very easy questions based on the content. make the questions as easy as possible."

    elif difficulty == 2:
        temperature = 0.2
        prompt = f"I have read a text about the following topic: {summary}. Please generate 5 easy questions based on the content. make the questions with a below average difficulty."

    elif difficulty == 3:
        temperature = 0.3
        prompt = f"I have read a text about the following topic: {summary}. Please generate 5 average questions based on the content. make the questions as average as possible."

    elif difficulty == 4:
        temperature = 0.4
        prompt = f"I have read a text about the following topic: {summary}. Please generate 5 pretty hard questions based on the content. make the questions pretty hard to understand, but not too much."

    elif difficulty == 5:
        temperature = 0.5
        prompt = f"I have read a text about the following topic: {summary}. Please generate 5 very hard questions based on the content. make the questions as hard as possible."

    chat = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=temperature)

    # Template to use for the system message prompt
    template = """
        You are a helpful assistant that can generate questions based on the provided text: {docs}
        
        Your task is to generate 5 meaningful questions that relate to the content of the text.
        """

    system_message_prompt = SystemMessagePromptTemplate.from_template(template)

    # We wrap the system message into a ChatPromptTemplate
    chat_prompt = ChatPromptTemplate.from_messages([system_message_prompt])

    # Generate questions
    chain = LLMChain(llm=chat, prompt=chat_prompt)
    questions = chain.run(question=prompt, docs=summary)

    # Assume the model returns the questions separated by newlines, split them
    questions = questions.split('\n')

    return questions


def summarize_text(db, summary_detail):
    # Determine max_length based on the summary_detail
    if summary_detail == 1:
        summary_question = "Can you give me a simple, easy-to-understand summary of this text? Just cover the main points briefly."
    elif summary_detail == 2:
        summary_question = "Can you provide a summary of this text that includes important details but remains fairly concise?"
    elif summary_detail == 3:
        summary_question = "Can you summarize this text for me, covering all key points and some finer details? Aim for a balance between brevity and comprehensiveness."
    elif summary_detail == 4:
        summary_question = "Can you provide a detailed summary of this text, digging deeper into the main points and also touching on the less crucial aspects?"
    elif summary_detail == 5:
        summary_question = "Can you provide a very in-depth summary of this text, including as many details as possible? Leave nothing important out."

    summary, _ = get_response_from_query(
        db, summary_question, summary_needed=True)
    return summary


def get_response_from_query(db, query, summary_needed, k=4):
    """
    gpt-3.5-turbo can handle up to 4097 tokens. Setting the chunksize to 1000 and k to 4 maximizes
    the number of tokens to analyze.
    """

    docs = db.similarity_search(query, k=k)
    docs_page_content = " ".join([d.page_content for d in docs])

    chat = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=0.2)

    # Template to use for the system message prompt
    template = """
        You are a helpful assistant that can answer questions about youtube videos 
        based on the video's transcript: {docs}
        
        Only use the factual information from the transcript to answer the question.
        
        If you feel like you don't have enough information to answer the question, say "I don't know".
        
        Your answers should be verbose and detailed.
        """

    system_message_prompt = SystemMessagePromptTemplate.from_template(template)

    # Human question prompt
    if summary_needed == True:
        human_template = "Answer the following question: {question}"
    else:
        human_template = "Answer the following question, please keep it short and straight to the point: {question}"

    human_message_prompt = HumanMessagePromptTemplate.from_template(
        human_template)

    chat_prompt = ChatPromptTemplate.from_messages(
        [system_message_prompt, human_message_prompt]
    )

    chain = LLMChain(llm=chat, prompt=chat_prompt)

    response = chain.run(question=query, docs=docs_page_content)
    response = response.replace("\n", "")
    return response, docs

def example_questions(summary):
    prompt = f"I have read a text about the following topic: {summary}. Please generate 3 simple example questions a student could ask for clarification about the topic based on the content."
    
    chat = ChatOpenAI(model_name="gpt-3.5-turbo")

    # Template to use for the system message prompt
    template = """
        You are a helpful assistant that can generate questions based on the provided text: {docs}
        
        Your task is to generate 3 short example questions that relate to the content of the text.
        """

    system_message_prompt = SystemMessagePromptTemplate.from_template(template)

    # We wrap the system message into a ChatPromptTemplate
    chat_prompt = ChatPromptTemplate.from_messages([system_message_prompt])

    # Generate questions
    chain = LLMChain(llm=chat, prompt=chat_prompt)
    ex_questions = chain.run(question=prompt, docs=summary)

    # Assume the model returns the questions separated by newlines, split them
    ex_questions = ex_questions.split('\n')

    return ex_questions

def main():
    # Set up the main user interface
    st.set_page_config(page_title="Study Material Learner") 
    st.image("https://assets.website-files.com/5e318ddf83dd66053d55c38a/602ba7efd5e8b761ed988e2a_FBF%20logo%20wide.png",
             caption='Study Material Learner', use_column_width=True) # Puts the Feedback Fruits logo on the page
    
    # Sets title of the page
    st.title("Study Material Learner")

    # Shows a short introduction of the tool
    st.write("""
    This tool is designed to enhance the learning experience for students by providing interactive study resources. 
    It enables students to extract valuable information from different sources, such as PDF files and lectures, and transforms that information into engaging quizzes. 
    These quizzes are generated automatically based on the specific study materials provided by the student. By utilizing this tool, students can deepen their understanding of the subjects they are studying and foster their intellectual growth. 
    It empowers students to actively participate in their learning process and develop their cognitive abilities, ultimately leading to academic success and personal development.
    """)

    progress = st.progress(0) # Keeps track of progress of the program running
    knowledge_base = None  # Initialize knowledge_base

    file = st.file_uploader("**Upload your file**", type=[
                                 "pdf", "docx", "pptx"])

    if file is None:
        st.write("_Please enter your file_")

    if file is not None: # Check if a file has been uploaded
        # Determine the type of the file
        filetype = file.type.split('/')[1]

        if filetype == 'pdf':
            # Read the PDF file and extract its text content into the text variable
            pdf_reader = PdfReader(file)
            progress.progress(10)
            text = ""

            for page in pdf_reader.pages:
                text += page.extract_text()

            progress.progress(20)

            # Create a CharacterTextSplitter object to split the text into smaller chunks for processing efficiency
            text_splitter = CharacterTextSplitter(
                separator="\n", # Define the separator for splitting the text into chunks
                chunk_size=1000, # Set the size of each chunk to 1000 characters
                chunk_overlap=200, # Set the overlap between consecutive chunks to 200 characters
                length_function=len # Define a function to calculate the length of the text
            )
            chunks = text_splitter.split_text(text) # Split the text into smaller chunks using the defined text splitter

            progress.progress(30)

            # Create an OpenAIEmbeddings object to generate embeddings for the text chunks
            embeddings = OpenAIEmbeddings()

            # Create a knowledge base using FAISS indexing from the text chunks and their embeddings
            knowledge_base = FAISS.from_texts(chunks, embeddings)

            progress.progress(40)
            
            # Set up a summary expander to generate a summary of the PDF
            # and allow users to interactively adjust the summary length
            summary_expander = st.expander('Summary of file')
            with summary_expander:
                st.session_state.summary_detail = None # Initialize summary detail
            
                if 'last_summary_detail' not in st.session_state:
                    st.session_state.last_summary_detail = None # Initialize the previous summary detail

                # Slider to input how long the summary should be
                st.session_state.summary_detail = st.slider(
                    'Set the lenght of your summary', 1, 5, 3)
                
                generate_summary = st.button("Generate Summary") # Generate the summary button

                # Check if `generate_summary` is True or if the `summary_detail` in the session state is different from the last summary detail
                if generate_summary or st.session_state.summary_detail != st.session_state.last_summary_detail:
                    st.session_state.last_summary_detail = st.session_state.summary_detail
                    # Update the last summary detail level if it has changed
                    # If a knowledge base exists, summarize the text using the given summary detail level
                    if knowledge_base:
                        summary = summarize_text(
                            knowledge_base, st.session_state.summary_detail) # Use the given summary detail level to summarize the text
                        st.session_state.summary = summary # Assigns the generated summary to the summary key in the session state
                        if generate_summary:
                            st.write(summary) # Write the summary
                        progress.progress(70)
                    else:
                        st.write("No summary available")

            # Set up a chat expander to allow users to interact with the document through questions
            chat_expander = st.expander('Chat with your file')
            with chat_expander:
                # Check if 'conversation_history' key is present in the session state
                if 'conversation_history' not in st.session_state:
                    st.session_state.conversation_history = [] # If 'conversation_history' key is not present, initialize it as an empty list

                # Check if 'user_question' key is present in the session state
                if 'user_question' not in st.session_state:
                    st.session_state.user_question = "" # If 'user_question' key is not present, initialize it as an empty string

                # Create a placeholder for the conversation history
                conversation_history_placeholder = st.empty()

                # Generate example questions
                if 'ex_questions' not in st.session_state: # If example questions are not available/present in the session state, generate them
                    ex_questions = example_questions(st.session_state.summary)
                    st.session_state.ex_questions = ex_questions # Store the generated example questions in the session state
                    
                else: # If example questions are already available in the session state, retrieve them
                    ex_questions = st.session_state.ex_questions                
                
                with st.form(key='chat_form'):
                    st.session_state.user_question = st.text_input("Ask a question about your PDF:", value=st.session_state.user_question)
                    # Where the user can write a question
                    selected_question = st.selectbox(" ", ['What kind of questions could I ask?'] + st.session_state.ex_questions)
                    # Where the user can select one of the example questions

                    submit_button = st.form_submit_button('Ask')
                    if submit_button:
                        if st.session_state.user_question: # If the user typed in a question use that input as the question to be processed
                            question = st.session_state.user_question
                        elif selected_question: # If the user selected a question use that input as the question to be processed
                            question = selected_question
                        else: # If the user didnt type or select a question there is no question
                            question = None

                        if question: # If there is a question begin further processing
                            st.session_state.conversation_history.append(('You', question)) # Append the question to the conversation history

                            docs = knowledge_base.similarity_search(question) # Perform similarity search on the knowledge base using the question
                            llm = OpenAI()
                            chain = load_qa_chain(llm, chain_type="stuff") # Load the question-answering chain model

                            with get_openai_callback() as cb: # Set up a callback to interact with OpenAI API
                                response = chain.run(input_documents=docs, question=question) # Generate a response to the question using the chain model

                            try:
                                st.session_state.conversation_history.append(('Bot', response)) # Append the response to the conversation history
                            except Exception as e:
                                st.error(f"An error occurred: {e}") # Display an error message if there was an exception

                            st.session_state.user_question = "" # Clear the user's question from the input field
                            selected_question = "" # Clear the selected example question
                            st.session_state.ex_questions = ex_questions

                        else:
                            st.warning("Please enter a question or select an example question.") # Display a warning message if no question was provided

                    # Display the conversation history
                    conversation_html = ""
                    if st.session_state.conversation_history: # Check if conversation history exists
                        for role, message in st.session_state.conversation_history: # Iterate over each role and message in the conversation history
                            if role == 'You': # If role is the user
                                conversation_html += f"<p><b>{role}:</b> {message}</p>" # Format the message as 'You: [message]'
                            else: # if role is the bot
                                conversation_html += f"<p><b>🤖:</b> {message}</p>" # Format the message as '🤖: [message]'
                    
                    # Render the conversation history in the app, displaying the messages exchanged between the user and the bot
                    conversation_history_placeholder.markdown(conversation_html, unsafe_allow_html=True)


            # Set up a questions expander to generate quiz questions based on the PDF content
            questions_expander = st.expander('Generated exam-like questions')
            with questions_expander:
                progress.progress(100)

                st.session_state.difficulty = None # Initialize the difficulty

                if 'last_difficulty' not in st.session_state: # Initialize the previous difficulty
                    st.session_state.last_difficulty = None

                # Difficult slider
                st.session_state.difficulty = st.slider(
                    'Set the difficulty level of your quiz questions', 1, 5, 3)

                if 'questions_generated' not in st.session_state: # Check if the flag for generated questions is not stored in the session state
                    st.session_state.questions_generated = False # Initialize the flag as False

                button_label = "Adjust Difficulty of Questions" if st.session_state.questions_generated else "Generate Questions"
                # If questions have been generated, set the button label as "Adjust Difficulty of Questions",
                # otherwise set it as "Generate Questions"
                
                generate_questions_button = st.button(button_label)
                # Create a button with the dynamically determined label ("Adjust Difficulty of Questions" or "Generate Questions")

                if generate_questions_button: # If the "Generate Questions" button is clicked
                    st.session_state.last_difficulty = st.session_state.difficulty # Store the current difficulty level in the session state as the last used difficulty level
                    questions = generate_questions(
                        st.session_state.summary, st.session_state.difficulty) # Generate questions based on the summary and the current difficulty level
                    st.session_state.questions = questions # Store the generated questions in the session state
                    st.session_state.answers = {}  # Reset the state for storing user answers
                    st.session_state.checked_answers = {}  # Reset the state for storing checked answers
                    st.session_state.questions_generated = True # Set the flag indicating that questions have been generated to True              

                # If the state for storing answers is not present in the session state, initialize it as an empty dictionary
                if 'answers' not in st.session_state:
                    st.session_state.answers = {}

                # If the state for storing questions is not present in the session state, initialize it as an empty list
                if 'questions' not in st.session_state:
                    st.session_state.questions = []

                if st.session_state.questions:  # Check if questions exist
                    st.write('Generated Questions') # Display a heading indicating that questions have been generated

                    for idx, question in enumerate(st.session_state.questions, start=1):
                        st.write(f"Q{idx}: {question}") # Display the question with its corresponding index
                        user_answer = st.text_input(
                            f"Your Answer for Q{idx}") # Create a text input field for the user to enter their answer for the question
                        st.session_state.answers[idx] = {
                            'question': question, 'answer': user_answer} # Store the user's answer for the question in the answers state

                        check_answere_button_placeholder = st.empty()
                        check_answere_button = check_answere_button_placeholder.button(f'Check Answer {idx}') # Button to check the answer

                        # Wait until the check answer button is clicked before proceeding 
                        if check_answere_button == False:
                            while not check_answere_button:
                                pass

                        check_answere_button_placeholder.empty()

                        # Check if the user's answer is correct by comparing it with the correct generated answer from the OpenAI model
                        is_correct, correct_answer = check_answer(
                            question, user_answer, st.session_state.summary, embeddings)
                        
                        if is_correct: # Tells the user that the answer is correct
                            st.success("Your answer is correct!")
                            st.markdown("---")  # Horizontal line between question answers

                        if is_correct == False: # Tells the user that the answer is correct
                            st.markdown("Your answer is **not** correct. The correct answer should be:")
                            st.markdown(correct_answer)
                            st.markdown("---")  # Horizontal line between question answers
                        
                        # If the state for storing checked answers is not present in the session state, initialize it as an empty dictionary
                        if 'checked_answers' not in st.session_state:
                            st.session_state.checked_answers = {}

                        # Store the information about the checked answer (whether it is correct or not) in the checked answers state
                        st.session_state.checked_answers[idx] = {
                            'is_correct': is_correct, 'correct_answer': correct_answer}
                    
                        next_question_button_placeholder = st.empty() # Button which can be clicked when the user wants to see the next question
                        # If the current question is not the last (6th) question, create a button labeled as "Question {idx + 1}" using the placeholder
                        if idx+1 is not 6:
                            next_question_button = next_question_button_placeholder.button(f'Question {idx + 1}') 

                        # If the current question is the last (6th) question, create a button labeled as "Check all answers" using the placeholder
                        if idx+1 is 6:
                            next_question_button = next_question_button_placeholder.button('Check all answers')

                        # Wait until the next question button is clicked (False is returned) before proceeding to the next iteration of the loop
                        if next_question_button == False:
                            while not next_question_button:
                                pass
                        
                        next_question_button_placeholder.empty()

                                       
                    # Calculate and display the user's score and progress only if all answers have been checked.
                    # It checks if 'checked_answers' is present in the session state and if the number of checked answers is equal to the number of questions.
                    if 'checked_answers' in st.session_state and len(st.session_state.checked_answers) == len(st.session_state.questions):
                        # Calculates the score by counting the number of correct answers in the checked answers.
                        # The score is displayed along with the total number of questions.
                        score = len(
                            [data for data in st.session_state.checked_answers.values() if data['is_correct']])
                        st.write(
                            f"Your score: {score} out of {len(st.session_state.questions)}")
                        # A progress bar is also created and its progress is set based on the score.
                        progress_bar = st.progress(0)
                        progress_percentage = score / \
                            len(st.session_state.questions)
                        progress_bar.progress(progress_percentage)

        elif filetype == 'vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc = Document(file) # Create a Document object to read the content from the file
            progress.progress(10)
            text = "" # Initialize an empty string to store the text content of the document

            # Extract the text from each paragraph in the document and concatenate them into a single text
            for para in doc.paragraphs:
                text += para.text

            progress.progress(20)

            # Create a CharacterTextSplitter object to split the text into smaller chunks for processing efficiency
            text_splitter = CharacterTextSplitter(
                separator="\n", # Define the separator for splitting the text into chunks
                chunk_size=1000, # Set the size of each chunk to 1000 characters
                chunk_overlap=200, # Set the overlap between consecutive chunks to 200 characters
                length_function=len # Define a function to calculate the length of the text
            )
            chunks = text_splitter.split_text(text) # Split the text into smaller chunks using the defined text splitter

            progress.progress(30)

            # Create an OpenAIEmbeddings object to generate embeddings for the text chunks
            embeddings = OpenAIEmbeddings()

            # Create a knowledge base using FAISS indexing from the text chunks and their embeddings
            knowledge_base = FAISS.from_texts(chunks, embeddings)

            progress.progress(40)
            
            # Set up a summary expander to generate a summary of the PDF
            # and allow users to interactively adjust the summary length
            summary_expander = st.expander('Summary of file')
            with summary_expander:
                st.session_state.summary_detail = None # Initialize summary detail

                if 'last_summary_detail' not in st.session_state:
                    st.session_state.last_summary_detail = None  # Initialize the previous summary detail

                # Slider to input how long the summary should be
                st.session_state.summary_detail = st.slider(
                    'Set the lenght of your summary', 1, 5, 3)
                
                generate_summary = st.button("Generate Summary") # Generate the summary button

                # Check if `generate_summary` is True or if the `summary_detail` in the session state is different from the last summary detail
                if generate_summary or st.session_state.summary_detail != st.session_state.last_summary_detail:
                    st.session_state.last_summary_detail = st.session_state.summary_detail
                    # Update the last summary detail level if it has changed
                    # If a knowledge base exists, summarize the text using the given summary detail level
                    if knowledge_base:
                        summary = summarize_text(
                            knowledge_base, st.session_state.summary_detail) # Use the given summary detail level to summarize the text
                        st.session_state.summary = summary # Assigns the generated summary to the summary key in the session state
                        if generate_summary:
                            st.write(summary) # Write the summary
                        progress.progress(70)
                    else:
                        st.write("No summary available")

            # Set up a chat expander to allow users to interact with the document through questions
            chat_expander = st.expander('Chat with your file')
            with chat_expander:
                # Check if 'conversation_history' key is present in the session state
                if 'conversation_history' not in st.session_state:
                    st.session_state.conversation_history = [] # If 'conversation_history' key is not present, initialize it as an empty list

                # Check if 'user_question' key is present in the session state
                if 'user_question' not in st.session_state:
                    st.session_state.user_question = "" # If 'user_question' key is not present, initialize it as an empty string

                # Create a placeholder for the conversation history
                conversation_history_placeholder = st.empty()

                # Generate example questions
                if 'ex_questions' not in st.session_state: # If example questions are not available/present in the session state, generate them
                    ex_questions = example_questions(st.session_state.summary)
                    st.session_state.ex_questions = ex_questions # Store the generated example questions in the session state
                    
                else: # If example questions are already available in the session state, retrieve them
                    ex_questions = st.session_state.ex_questions                

                with st.form(key='chat_form'):
                    st.session_state.user_question = st.text_input("Ask a question about your PDF:", value=st.session_state.user_question)
                    # Where the user can write a question
                    selected_question = st.selectbox(" ", ['What kind of questions could I ask?'] + st.session_state.ex_questions)
                    # Where the user can select one of the example questions

                    submit_button = st.form_submit_button('Ask')
                    if submit_button:
                        if st.session_state.user_question: # If the user typed in a question use that input as the question to be processed
                            question = st.session_state.user_question
                        elif selected_question: # If the user selected a question use that input as the question to be processed
                            question = selected_question
                        else: # If the user didnt type or select a question there is no question
                            question = None

                        if question: # If there is a question begin further processing
                            st.session_state.conversation_history.append(('You', question)) # Append the question to the conversation history

                            docs = knowledge_base.similarity_search(question) # Perform similarity search on the knowledge base using the question
                            llm = OpenAI()
                            chain = load_qa_chain(llm, chain_type="stuff") # Load the question-answering chain model

                            with get_openai_callback() as cb: # Set up a callback to interact with OpenAI API
                                response = chain.run(input_documents=docs, question=question) # Generate a response to the question using the chain model

                            try:
                                st.session_state.conversation_history.append(('Bot', response)) # Append the response to the conversation history
                            except Exception as e:
                                st.error(f"An error occurred: {e}") # Display an error message if there was an exception

                            st.session_state.user_question = "" # Clear the user's question from the input field
                            selected_question = "" # Clear the selected example question
                            st.session_state.ex_questions = ex_questions

                        else:
                            st.warning("Please enter a question or select an example question.") # Display a warning message if no question was provided

                    # Display the conversation history
                    conversation_html = ""
                    if st.session_state.conversation_history: # Check if conversation history exists
                        for role, message in st.session_state.conversation_history: # Iterate over each role and message in the conversation history
                            if role == 'You': # If role is the user
                                conversation_html += f"<p><b>{role}:</b> {message}</p>" # Format the message as 'You: [message]'
                            else: # if role is the bot
                                conversation_html += f"<p><b>🤖:</b> {message}</p>" # Format the message as '🤖: [message]'
                    
                    # Render the conversation history in the app, displaying the messages exchanged between the user and the bot
                    conversation_history_placeholder.markdown(conversation_html, unsafe_allow_html=True)
            
            # Set up a questions expander to generate quiz questions based on the PDF content
            questions_expander = st.expander('Generated exam-like questions')
            with questions_expander:
                progress.progress(100)

                st.session_state.difficulty = None # Initialize the difficulty

                if 'last_difficulty' not in st.session_state: # Initialize the previous difficulty
                    st.session_state.last_difficulty = None

                # Difficult slider
                st.session_state.difficulty = st.slider(
                    'Set the difficulty level of your quiz questions', 1, 5, 3)

                if 'questions_generated' not in st.session_state: # Check if the flag for generated questions is not stored in the session state
                    st.session_state.questions_generated = False # Initialize the flag as False

                button_label = "Adjust Difficulty of Questions" if st.session_state.questions_generated else "Generate Questions"
                # If questions have been generated, set the button label as "Adjust Difficulty of Questions",
                # otherwise set it as "Generate Questions"

                generate_questions_button = st.button(button_label)
                # Create a button with the dynamically determined label ("Adjust Difficulty of Questions" or "Generate Questions")

                if generate_questions_button: # If the "Generate Questions" button is clicked
                    st.session_state.last_difficulty = st.session_state.difficulty # Store the current difficulty level in the session state as the last used difficulty level
                    questions = generate_questions(
                        st.session_state.summary, st.session_state.difficulty) # Generate questions based on the summary and the current difficulty level
                    st.session_state.questions = questions # Store the generated questions in the session state
                    st.session_state.answers = {}  # Reset the state for storing user answers
                    st.session_state.checked_answers = {}  # Reset the state for storing checked answers
                    st.session_state.questions_generated = True # Set the flag indicating that questions have been generated to True              

                # If the state for storing answers is not present in the session state, initialize it as an empty dictionary
                if 'answers' not in st.session_state:
                    st.session_state.answers = {}

                # If the state for storing questions is not present in the session state, initialize it as an empty list
                if 'questions' not in st.session_state:
                    st.session_state.questions = []

                if st.session_state.questions:  # Check if questions exist
                    st.write('Generated Questions') # Display a heading indicating that questions have been generated

                    for idx, question in enumerate(st.session_state.questions, start=1):
                        st.write(f"Q{idx}: {question}") # Display the question with its corresponding index
                        user_answer = st.text_input(
                            f"Your Answer for Q{idx}") # Create a text input field for the user to enter their answer for the question
                        st.session_state.answers[idx] = {
                            'question': question, 'answer': user_answer} # Store the user's answer for the question in the answers state

                        check_answere_button_placeholder = st.empty()
                        check_answere_button = check_answere_button_placeholder.button(f'Check Answer {idx}') # Button to check the answer

                        # Wait until the check answer button is clicked before proceeding 
                        if check_answere_button == False:
                            while not check_answere_button:
                                pass

                        check_answere_button_placeholder.empty()

                        # Check if the user's answer is correct by comparing it with the correct generated answer from the OpenAI model
                        is_correct, correct_answer = check_answer(
                            question, user_answer, st.session_state.summary, embeddings)
                        
                        if is_correct: # Tells the user that the answer is correct
                            st.success("Your answer is correct!")
                            st.markdown("---")  # Horizontal line between question answers

                        if is_correct == False: # Tells the user that the answer is correct
                            st.markdown("Your answer is **not** correct. The correct answer should be:")
                            st.markdown(correct_answer)
                            st.markdown("---")  # Horizontal line between question answers
                        
                        # If the state for storing checked answers is not present in the session state, initialize it as an empty dictionary
                        if 'checked_answers' not in st.session_state:
                            st.session_state.checked_answers = {}

                        # Store the information about the checked answer (whether it is correct or not) in the checked answers state
                        st.session_state.checked_answers[idx] = {
                            'is_correct': is_correct, 'correct_answer': correct_answer}
                    
                        next_question_button_placeholder = st.empty() # Button which can be clicked when the user wants to see the next question
                        # If the current question is not the last (6th) question, create a button labeled as "Question {idx + 1}" using the placeholder
                        if idx+1 is not 6:
                            next_question_button = next_question_button_placeholder.button(f'Question {idx + 1}') 

                        # If the current question is the last (6th) question, create a button labeled as "Check all answers" using the placeholder
                        if idx+1 is 6:
                            next_question_button = next_question_button_placeholder.button('Check all answers')

                        # Wait until the next question button is clicked (False is returned) before proceeding to the next iteration of the loop
                        if next_question_button == False:
                            while not next_question_button:
                                pass
                        
                        next_question_button_placeholder.empty()
                       
                    # Calculate and display the user's score and progress only if all answers have been checked.
                    # It checks if 'checked_answers' is present in the session state and if the number of checked answers is equal to the number of questions.
                    if 'checked_answers' in st.session_state and len(st.session_state.checked_answers) == len(st.session_state.questions):
                        # Calculates the score by counting the number of correct answers in the checked answers.
                        # The score is displayed along with the total number of questions.
                        score = len(
                            [data for data in st.session_state.checked_answers.values() if data['is_correct']])
                        st.write(
                            f"Your score: {score} out of {len(st.session_state.questions)}")
                        # A progress bar is also created and its progress is set based on the score.
                        progress_bar = st.progress(0)
                        progress_percentage = score / \
                            len(st.session_state.questions)
                        progress_bar.progress(progress_percentage)

        elif filetype == 'vnd.openxmlformats-officedocument.presentationml.presentation':
            text = extract_text_from_pptx(file) # Extract text from pptx file
            
            progress.progress(20)

            # Create a CharacterTextSplitter object to split the text into smaller chunks for processing efficiency
            text_splitter = CharacterTextSplitter(
                separator="\n", # Define the separator for splitting the text into chunks
                chunk_size=1000, # Set the size of each chunk to 1000 characters
                chunk_overlap=200, # Set the overlap between consecutive chunks to 200 characters
                length_function=len # Define a function to calculate the length of the text
            )
            chunks = text_splitter.split_text(text) # Split the text into smaller chunks using the defined text splitter

            progress.progress(30)

            # Create an OpenAIEmbeddings object to generate embeddings for the text chunks
            embeddings = OpenAIEmbeddings()

            # Create a knowledge base using FAISS indexing from the text chunks and their embeddings
            knowledge_base = FAISS.from_texts(chunks, embeddings)

            progress.progress(40)

            # Set up a summary expander to generate a summary of the presentation
            # and allow users to interactively adjust the summary length
            summary_expander = st.expander('Summary of file')
            with summary_expander:
                st.session_state.summary_detail = None # Initialize summary detail
            
                if 'last_summary_detail' not in st.session_state:
                    st.session_state.last_summary_detail = None # Initialize the previous summary detail

                # Slider to input how long the summary should be
                st.session_state.summary_detail = st.slider(
                    'Set the lenght of your summary', 1, 5, 3)
                
                generate_summary = st.button("Generate Summary") # Generate the summary button

                # Check if `generate_summary` is True or if the `summary_detail` in the session state is different from the last summary detail
                if generate_summary or st.session_state.summary_detail != st.session_state.last_summary_detail:
                    st.session_state.last_summary_detail = st.session_state.summary_detail
                    # Update the last summary detail level if it has changed
                    # If a knowledge base exists, summarize the text using the given summary detail level
                    if knowledge_base:
                        summary = summarize_text(
                            knowledge_base, st.session_state.summary_detail) # Use the given summary detail level to summarize the text
                        st.session_state.summary = summary # Assigns the generated summary to the summary key in the session state
                        if generate_summary:
                            st.write(summary) # Write the summary
                        progress.progress(70)
                    else:
                        st.write("No summary available")

            # Set up a chat expander to allow users to interact with the presentation through questions
            chat_expander = st.expander('Chat with your file')
            with chat_expander:
                # Check if 'conversation_history' key is present in the session state
                if 'conversation_history' not in st.session_state:
                    st.session_state.conversation_history = [] # If 'conversation_history' key is not present, initialize it as an empty list

                # Check if 'user_question' key is present in the session state
                if 'user_question' not in st.session_state:
                    st.session_state.user_question = "" # If 'user_question' key is not present, initialize it as an empty string

                # Create a placeholder for the conversation history
                conversation_history_placeholder = st.empty()

                # Generate example questions
                if 'ex_questions' not in st.session_state: # If example questions are not available/present in the session state, generate them
                    ex_questions = example_questions(st.session_state.summary)
                    st.session_state.ex_questions = ex_questions # Store the generated example questions in the session state
                    
                else: # If example questions are already available in the session state, retrieve them
                    ex_questions = st.session_state.ex_questions                
                
                with st.form(key='chat_form'):
                    st.session_state.user_question = st.text_input("Ask a question about your presentation:", value=st.session_state.user_question)
                    # Where the user can write a question
                    selected_question = st.selectbox(" ", ['What kind of questions could I ask?'] + st.session_state.ex_questions)
                    # Where the user can select one of the example questions

                    submit_button = st.form_submit_button('Ask')
                    if submit_button:
                        if st.session_state.user_question: # If the user typed in a question use that input as the question to be processed
                            question = st.session_state.user_question
                        elif selected_question: # If the user selected a question use that input as the question to be processed
                            question = selected_question
                        else: # If the user didnt type or select a question there is no question
                            question = None

                        if question: # If there is a question begin further processing
                            st.session_state.conversation_history.append(('You', question)) # Append the question to the conversation history

                            docs = knowledge_base.similarity_search(question) # Perform similarity search on the knowledge base using the question
                            llm = OpenAI()
                            chain = load_qa_chain(llm, chain_type="stuff") # Load the question-answering chain model

                            with get_openai_callback() as cb: # Set up a callback to interact with OpenAI API
                                response = chain.run(input_documents=docs, question=question) # Generate a response to the question using the chain model

                            try:
                                st.session_state.conversation_history.append(('Bot', response)) # Append the response to the conversation history
                            except Exception as e:
                                st.error(f"An error occurred: {e}") # Display an error message if there was an exception

                            st.session_state.user_question = "" # Clear the user's question from the input field
                            selected_question = "" # Clear the selected example question
                            st.session_state.ex_questions = ex_questions

                        else:
                            st.warning("Please enter a question or select an example question.") # Display a warning message if no question was provided

                    # Display the conversation history
                    conversation_html = ""
                    if st.session_state.conversation_history: # Check if conversation history exists
                        for role, message in st.session_state.conversation_history: # Iterate over each role and message in the conversation history
                            if role == 'You': # If role is the user
                                conversation_html += f"<p><b>{role}:</b> {message}</p>" # Format the message as 'You: [message]'
                            else: # If role is the bot
                                conversation_html += f"<p><b>🤖:</b> {message}</p>" # Format the message as '🤖: [message]'
                    
                    # Render the conversation history in the app, displaying the messages exchanged between the user and the bot
                    conversation_history_placeholder.markdown(conversation_html, unsafe_allow_html=True)

            # Set up a questions expander to generate quiz questions based on the file content
            questions_expander = st.expander('Generated exam-like questions')
            with questions_expander:
                progress.progress(100)

                st.session_state.difficulty = None # Initialize the difficulty

                if 'last_difficulty' not in st.session_state: # Initialize the previous difficulty
                    st.session_state.last_difficulty = None

                # Difficult slider
                st.session_state.difficulty = st.slider(
                    'Set the difficulty level of your quiz questions', 1, 5, 3)

                if 'questions_generated' not in st.session_state: # Check if the flag for generated questions is not stored in the session state
                    st.session_state.questions_generated = False # Initialize the flag as False

                button_label = "Adjust Difficulty of Questions" if st.session_state.questions_generated else "Generate Questions"
                # If questions have been generated, set the button label as "Adjust Difficulty of Questions",
                # otherwise set it as "Generate Questions"
                
                generate_questions_button = st.button(button_label)
                # Create a button with the dynamically determined label ("Adjust Difficulty of Questions" or "Generate Questions")

                if generate_questions_button: # If the "Generate Questions" button is clicked
                    st.session_state.last_difficulty = st.session_state.difficulty # Store the current difficulty level in the session state as the last used difficulty level
                    questions = generate_questions(
                        st.session_state.summary, st.session_state.difficulty) # Generate questions based on the summary and the current difficulty level
                    st.session_state.questions = questions # Store the generated questions in the session state
                    st.session_state.answers = {}  # Reset the state for storing user answers
                    st.session_state.checked_answers = {}  # Reset the state for storing checked answers
                    st.session_state.questions_generated = True # Set the flag indicating that questions have been generated to True              

                # If the state for storing answers is not present in the session state, initialize it as an empty dictionary
                if 'answers' not in st.session_state:
                    st.session_state.answers = {}

                # If the state for storing questions is not present in the session state, initialize it as an empty list
                if 'questions' not in st.session_state:
                    st.session_state.questions = []

                if st.session_state.questions:  # Check if questions exist
                    st.write('Generated Questions') # Display a heading indicating that questions have been generated
                    
                    for idx, question in enumerate(st.session_state.questions, start=1):
                        st.write(f"Q{idx}: {question}") # Display the question with its corresponding index
                        user_answer = st.text_input(
                            f"Your Answer for Q{idx}") # Create a text input field for the user to enter their answer for the question
                        st.session_state.answers[idx] = {
                            'question': question, 'answer': user_answer} # Store the user's answer for the question in the answers state

                        check_answere_button_placeholder = st.empty()
                        check_answere_button = check_answere_button_placeholder.button(f'Check Answer {idx}') # Button to check the answer

                        # Wait until the check answer button is clicked before proceeding 
                        if check_answere_button == False:
                            while not check_answere_button:
                                pass

                        check_answere_button_placeholder.empty()

                        # Check if the user's answer is correct by comparing it with the correct generated answer from the OpenAI model
                        is_correct, correct_answer = check_answer(
                            question, user_answer, st.session_state.summary, embeddings)
                        
                        if is_correct: # Tells the user that the answer is correct
                            st.success("Your answer is correct!")
                            st.markdown("---")  # Horizontal line between question answers

                        if is_correct == False: # Tells the user that the answer is correct
                            st.markdown("Your answer is **not** correct. The correct answer should be:")
                            st.markdown(correct_answer)
                            st.markdown("---")  # Horizontal line between question answers
                        
                        # If the state for storing checked answers is not present in the session state, initialize it as an empty dictionary
                        if 'checked_answers' not in st.session_state:
                            st.session_state.checked_answers = {}

                        # Store the information about the checked answer (whether it is correct or not) in the checked answers state
                        st.session_state.checked_answers[idx] = {
                            'is_correct': is_correct, 'correct_answer': correct_answer}
                    
                        next_question_button_placeholder = st.empty() # Button which can be clicked when the user wants to see the next question
                        # If the current question is not the last (6th) question, create a button labeled as "Question {idx + 1}" using the placeholder
                        if idx+1 is not 6:
                            next_question_button = next_question_button_placeholder.button(f'Question {idx + 1}') 

                        # If the current question is the last (6th) question, create a button labeled as "Check all answers" using the placeholder
                        if idx+1 is 6:
                            next_question_button = next_question_button_placeholder.button('Check all answers')

                        # Wait until the next question button is clicked (False is returned) before proceeding to the next iteration of the loop
                        if next_question_button == False:
                            while not next_question_button:
                                pass
                        
                        next_question_button_placeholder.empty()

                                       
                    # Calculate and display the user's score and progress only if all answers have been checked.
                    # It checks if 'checked_answers' is present in the session state and if the number of checked answers is equal to the number of questions.
                    if 'checked_answers' in st.session_state and len(st.session_state.checked_answers) == len(st.session_state.questions):
                        # Calculates the score by counting the number of correct answers in the checked answers.
                        # The score is displayed along with the total number of questions.
                        score = len(
                            [data for data in st.session_state.checked_answers.values() if data['is_correct']])
                        st.write(
                            f"Your score: {score} out of {len(st.session_state.questions)}")
                        # A progress bar is also created and its progress is set based on the score.
                        progress_bar = st.progress(0)
                        progress_percentage = score / \
                            len(st.session_state.questions)
                        progress_bar.progress(progress_percentage)

if __name__ == '__main__':
    main()
